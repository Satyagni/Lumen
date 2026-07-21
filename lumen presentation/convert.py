import os
import asyncio
from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.xmlchemy import OxmlElement

async def capture_slides():
    html_path = os.path.abspath("Lumen_Presentation.html")
    file_url = f"file:///{html_path.replace(os.sep, '/')}"
    print(f"Opening {file_url}...")
    
    screenshot_dir = "temp_screenshots"
    os.makedirs(screenshot_dir, exist_ok=True)
    
    async with async_playwright() as p:
        # Launch browser in headless mode
        browser = await p.chromium.launch(headless=True)
        context = await browser.new_context(viewport={"width": 1920, "height": 1080})
        page = await context.new_page()
        
        # Navigate to local file
        await page.goto(file_url)
        await page.wait_for_load_state("networkidle")
        
        # Hide presentation UI controls to make the PowerPoint slides look clean and native
        await page.evaluate("""() => {
            const p = document.getElementById('progress-bar');
            if (p) p.style.display = 'none';
            const n = document.getElementById('nav-dots');
            if (n) n.style.display = 'none';
            const c = document.getElementById('controls');
            if (c) c.style.display = 'none';
        }""")
        
        # Get total slides
        total_slides = await page.evaluate("document.querySelectorAll('.slide').length")
        print(f"Detected {total_slides} slides.")
        
        captured_slides = [] # Elements will be (slide_index, step_index, is_new_slide, file_path)
        
        for i in range(total_slides):
            print(f"Analyzing slide {i+1}/{total_slides} for stagger animations...")
            
            # Find how many stagger levels exist on this slide
            num_steps = await page.evaluate("""(slideIndex) => {
                const order = ['stagger-scale', 'stagger-1', 'stagger-2', 'stagger-3', 'stagger-4', 'stagger-5', 'stagger-6'];
                const slide = document.getElementById(`slide-${slideIndex}`);
                if (!slide) return 0;
                const staggers = slide.querySelectorAll('.stagger-1, .stagger-2, .stagger-3, .stagger-4, .stagger-5, .stagger-6, .stagger-scale');
                const levels = new Set();
                staggers.forEach(el => {
                    el.classList.forEach(cls => {
                        if (cls.startsWith('stagger-')) {
                            levels.add(cls);
                        }
                    });
                });
                const sortedLevels = order.filter(lvl => levels.has(lvl));
                return sortedLevels.length;
            }""", i)
            
            # Total steps to capture (Step 0: base content only; Step 1..N: show progressive layers)
            steps_to_capture = max(1, num_steps + 1)
            print(f"  Slide {i+1} has {num_steps} stagger levels. Capturing {steps_to_capture} steps.")
            
            for step in range(steps_to_capture):
                print(f"    Capturing step {step+1}/{steps_to_capture}...")
                
                await page.evaluate("""([slideIndex, step]) => {
                    // Activate slide i and deactivate others
                    document.querySelectorAll('.slide').forEach((s, idx) => {
                        s.classList.remove('enter-from-right', 'enter-from-left', 'exit-to-left', 'exit-to-right');
                        s.classList.toggle('active', idx === slideIndex);
                    });
                    if (typeof updateUI === 'function') {
                        updateUI(slideIndex);
                    }
                    
                    // Adjust stagger visibilities manually for the step
                    const order = ['stagger-scale', 'stagger-1', 'stagger-2', 'stagger-3', 'stagger-4', 'stagger-5', 'stagger-6'];
                    const slide = document.getElementById(`slide-${slideIndex}`);
                    if (!slide) return;
                    
                    const staggers = slide.querySelectorAll('.stagger-1, .stagger-2, .stagger-3, .stagger-4, .stagger-5, .stagger-6, .stagger-scale');
                    
                    // Find unique levels present on this slide
                    const levels = new Set();
                    staggers.forEach(el => {
                        el.classList.forEach(cls => {
                            if (cls.startsWith('stagger-')) {
                                levels.add(cls);
                            }
                        });
                    });
                    const sortedLevels = order.filter(lvl => levels.has(lvl));
                    
                    staggers.forEach(el => {
                        // Disable CSS animations/transitions for instant render control
                        el.style.animation = 'none';
                        el.style.transition = 'none';
                        
                        let shouldShow = true;
                        sortedLevels.forEach((lvl, stepIdx) => {
                            if (el.classList.contains(lvl) && stepIdx >= step) {
                                shouldShow = false;
                            }
                        });
                        
                        if (shouldShow) {
                            el.style.opacity = '1';
                            el.style.transform = 'none';
                        } else {
                            el.style.opacity = '0';
                            if (el.classList.contains('stagger-scale')) {
                                el.style.transform = 'scale(0.92)';
                            } else {
                                el.style.transform = 'translateY(16px)';
                            }
                        }
                    });
                }""", [i, step])
                
                # A brief pause to ensure browser rendering engine completes layout paint
                await asyncio.sleep(0.15)
                
                screenshot_path = os.path.join(screenshot_dir, f"slide_{i:03d}_step_{step:02d}.png")
                await page.screenshot(path=screenshot_path)
                
                is_new_slide = (step == 0)
                captured_slides.append((i, step, is_new_slide, screenshot_path))
                
        await browser.close()
    return captured_slides

def add_push_transition(slide, direction="l"):
    sld = slide._element
    # Create transition element
    transition = OxmlElement('p:transition')
    # p:push transition matches the HTML slider entering from the right
    push = OxmlElement('p:push')
    push.set('dir', direction)
    transition.append(push)
    sld.append(transition)

def create_presentation(captured_slides, output_path):
    print("Creating PPTX presentation with build transitions...")
    prs = Presentation()
    # Set slide dimensions to standard widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    for idx, (slide_idx, step_idx, is_new_slide, screenshot_path) in enumerate(captured_slides):
        slide = prs.slides.add_slide(blank_layout)
        # Add the screenshot image full-bleed
        slide.shapes.add_picture(screenshot_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        # Apply transition ONLY when moving to a brand new slide
        # (Step 0 of slides 1..N). We don't apply transitions on slide 0 step 0 or sub-steps.
        if is_new_slide and idx > 0:
            add_push_transition(slide, direction="l")
            
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}!")

async def main():
    captured_slides = await capture_slides()
    output_pptx = "Lumen_Presentation_Animated.pptx"
    create_presentation(captured_slides, output_pptx)
    
    # Clean up temp screenshots
    print("Cleaning up temporary screenshots...")
    for _, _, _, path in captured_slides:
        try:
            os.remove(path)
        except OSError:
            pass
    try:
        os.rmdir("temp_screenshots")
    except OSError:
        pass
    print("All tasks completed successfully!")

if __name__ == "__main__":
    from playwright.async_api import async_playwright
    asyncio.run(main())
